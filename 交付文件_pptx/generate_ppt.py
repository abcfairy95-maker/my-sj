#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate_ppt.py - Chuxiong Normal University RAG Intelligent Q&A Agent
Defense presentation PPT (10 slides, 16:9, purple/gold academic theme)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = r"C:\Users\Administrator\Desktop\workspace\交付文件_pptx"
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "答辩演示_完整版.pptx")

# Colors
PRIMARY         = RGBColor(0x6F, 0x4A, 0x8E)
PRIMARY_DARK    = RGBColor(0x4A, 0x2E, 0x66)
PRIMARY_LIGHT   = RGBColor(0x9B, 0x7E, 0xB7)
PRIMARY_LIGHTER = RGBColor(0xC4, 0xB3, 0xD8)
ACCENT          = RGBColor(0xD6, 0x9E, 0x2E)
ACCENT_LIGHT    = RGBColor(0xF0, 0xD0, 0x70)
BG_COLOR        = RGBColor(0xF7, 0xF4, 0xF9)
TEXT_DARK       = RGBColor(0x1E, 0x1F, 0x23)
TEXT_LIGHT      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY       = RGBColor(0x88, 0x88, 0x88)
CODE_BG         = RGBColor(0x2D, 0x2D, 0x2D)
CODE_TEXT_COLOR = RGBColor(0xE0, 0xE0, 0xE0)
GREEN           = RGBColor(0x27, 0xAE, 0x60)
RED             = RGBColor(0xE7, 0x4C, 0x3C)
CARD_BG         = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG         = RGBColor(0xF0, 0xEE, 0xF4)
BORDER_COLOR    = RGBColor(0xE0, 0xDC, 0xE8)

PURPLE_SHADES = [
    RGBColor(0x6F, 0x4A, 0x8E), RGBColor(0x7B, 0x5A, 0x9A),
    RGBColor(0x8A, 0x6A, 0xA6), RGBColor(0x9A, 0x7A, 0xB2),
    RGBColor(0xA8, 0x8A, 0xBE),
]

SLIDE_W = Inches(13.333); SLIDE_H = Inches(7.5)
CN_FONT = "Microsoft YaHei"; CODE_FONT = "Consolas"


# ============================================================
# HELPERS
# ============================================================

def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def _no_line(shape):
    shape.line.fill.background()

def _set_line(shape, color, w=Pt(0.5)):
    shape.line.color.rgb = color; shape.line.width = w

def add_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, st=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(st, l, t, w, h); _fill(s, color); _no_line(s); return s

def add_textbox(slide, l, t, w, h, text, fs=Pt(14), color=TEXT_DARK,
                bold=False, align=PP_ALIGN.LEFT, font=CN_FONT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = fs; p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tf

def add_rich_tf(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.paragraphs[0].text = ""
    return tf

def add_para(tf, text, fs=Pt(12), color=TEXT_DARK, bold=False,
             align=PP_ALIGN.LEFT, font=CN_FONT, sb=Pt(2), sa=Pt(2)):
    if len(tf.paragraphs) == 0 or tf.paragraphs[-1].text != "":
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[-1]
    p.text = text; p.font.size = fs; p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    p.space_before = sb; p.space_after = sa
    return p

def add_card(slide, l, t, w, h, color=CARD_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(shape, color); _set_line(shape, BORDER_COLOR, Pt(0.5))
    shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l + Pt(2), t + Pt(2), w, h)
    _fill(shadow, RGBColor(0xE8, 0xE4, 0xEE)); _no_line(shadow)
    sp_el = shadow._element; sp_el.getparent().remove(sp_el)
    shape._element.getparent().insert(
        list(shape._element.getparent()).index(shape._element), sp_el)
    return shape

def add_code_block(slide, l, t, w, h, code_text):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(bg, CODE_BG); _set_line(bg, RGBColor(0x44, 0x44, 0x44), Pt(0.5))
    pad = Pt(10)
    tb = slide.shapes.add_textbox(l + pad, t + pad, w - pad * 2, h - pad * 2)
    tf = tb.text_frame; tf.word_wrap = True; tf.paragraphs[0].text = ""
    for i, line in enumerate(code_text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line if line else " "; p.font.size = Pt(9)
        p.font.color.rgb = CODE_TEXT_COLOR; p.font.name = CODE_FONT
        p.space_before = Pt(0); p.space_after = Pt(0); p.line_spacing = Pt(13)
    return tf

def add_badge(slide, l, t, text, bg_color, text_color=TEXT_LIGHT, fs=Pt(10)):
    est_w = max(Inches(0.85), Pt(len(text) * 9 + 16))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, est_w, Pt(18))
    _fill(shape, bg_color); _no_line(shape); shape.adjustments[0] = 0.3
    tf = shape.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = text; p.font.size = fs
    p.font.color.rgb = text_color; p.font.bold = True
    p.font.name = CN_FONT; p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0); p.space_after = Pt(0)
    return shape

def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), PRIMARY)
    add_rect(slide, Inches(0), Inches(0.06), Inches(0.08), Inches(0.6), ACCENT)
    add_textbox(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.45),
                title, Pt(26), PRIMARY_DARK, True)
    add_rect(slide, Inches(0.4), Inches(0.62), Inches(12.5), Inches(0.015), BORDER_COLOR)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12), Inches(0.3),
                    subtitle, Pt(11), TEXT_GRAY)

def add_gold_arrow(slide, l, t):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, Inches(0.35), Inches(0.25))
    _fill(s, ACCENT); _no_line(s); return s

def fmt_cell(cell, text, fs=Pt(9), color=TEXT_DARK, bold=False, bg=None, align=PP_ALIGN.LEFT):
    cell.text = ""; p = cell.text_frame.paragraphs[0]; p.text = text
    p.font.size = fs; p.font.color.rgb = color; p.font.bold = bold
    p.font.name = CN_FONT; p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg


# ============================================================
# SLIDE 1 - COVER
# ============================================================
def create_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, PRIMARY_DARK)

    # circles top-right
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(0.6), Inches(2.2), Inches(2.2))
    _fill(c1, RGBColor(0x5A, 0x3E, 0x76)); _no_line(c1)
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.1), Inches(0.9), Inches(1.6), Inches(1.6))
    _fill(c2, PRIMARY); _no_line(c2)
    p = c2.text_frame.paragraphs[0]
    p.text = "RAG"; p.font.size = Pt(32); p.font.color.rgb = TEXT_LIGHT
    p.font.bold = True; p.font.name = CN_FONT; p.alignment = PP_ALIGN.CENTER

    # decorative dots
    for i, (dx, dy) in enumerate([(1.0, 6.6), (2.0, 6.8), (11.5, 6.0), (12.2, 6.4)]):
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(0.12), Inches(0.12))
        _fill(d, ACCENT if i % 2 == 0 else PRIMARY_LIGHT); _no_line(d)

    add_rect(slide, Inches(1.0), Inches(3.15), Inches(1.5), Inches(0.05), ACCENT)

    add_textbox(slide, Inches(1.0), Inches(1.8), Inches(10.5), Inches(1.3),
                "基于大模型与RAG的\n智能问答Agent搭建", Pt(38), TEXT_LIGHT, True)
    add_textbox(slide, Inches(1.0), Inches(3.4), Inches(10), Inches(0.55),
                "楚雄师范学院 - 校园智能问答助手", Pt(20), ACCENT_LIGHT)
    add_textbox(slide, Inches(1.0), Inches(4.2), Inches(10), Inches(0.4),
                "DeepSeek Chat  |  Qwen Embedding  |  ChromaDB  |  Gradio  |  Python",
                Pt(13), PRIMARY_LIGHTER)
    add_textbox(slide, Inches(1.0), Inches(5.0), Inches(5), Inches(0.35),
                "2025年 - 毕业设计答辩", Pt(13), PRIMARY_LIGHTER)
    add_textbox(slide, Inches(1.0), Inches(5.4), Inches(5), Inches(0.35),
                "计算机科学与技术学院", Pt(12), PRIMARY_LIGHT)


# ============================================================
# SLIDE 2 - PROJECT BACKGROUND
# ============================================================
def create_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "项目背景",
                  "大语言模型存在幻觉问题，RAG技术通过检索外部知识库有效缓解")

    # LEFT - Problem
    add_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.3))
    add_badge(slide, Inches(0.8), Inches(1.35), "X  问题：大模型幻觉", RED, TEXT_LIGHT, Pt(11))
    tf = add_rich_tf(slide, Inches(0.8), Inches(1.85), Inches(5.2), Inches(2.4))
    add_para(tf, "大语言模型（LLM）的幻觉问题：", Pt(12), TEXT_DARK, True)
    for item in [
        "LLM在缺乏训练数据时可能生成虚假、不准确的内容",
        "对于校园特定场景（宿舍、食堂、选课等），通用模型无法提供准确答案",
        "模型知识截止日期限制，无法获取最新校园政策与通知",
        "直接使用会导致用户体验差，信息可信度低",
    ]:
        add_para(tf, "  - " + item, Pt(11), TEXT_DARK)
    add_para(tf, "", Pt(4))
    add_para(tf, '  用户："楚雄师院图书馆几点关门？"', Pt(10), TEXT_GRAY)
    add_para(tf, '  LLM："我无法确定具体时间..."  <- 幻觉/无答案', Pt(10), RED)

    # RIGHT - Solution
    add_card(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(3.3))
    add_badge(slide, Inches(7.1), Inches(1.35), "OK  方案：RAG检索增强生成", GREEN, TEXT_LIGHT, Pt(11))
    tf = add_rich_tf(slide, Inches(7.1), Inches(1.85), Inches(5.2), Inches(2.4))
    add_para(tf, "RAG（Retrieval-Augmented Generation）：", Pt(12), TEXT_DARK, True)
    for item in [
        "先检索 -> 再增强 -> 后生成，将外部知识注入LLM",
        "将校园文档向量化存储，根据用户问题检索相关内容",
        "将检索到的文档片段作为Prompt上下文，引导LLM准确回答",
        "答案可溯源至具体文档，提升可信度",
    ]:
        add_para(tf, "  - " + item, Pt(11), TEXT_DARK)
    add_para(tf, "", Pt(4))
    add_para(tf, '  用户："楚雄师院图书馆几点关门？"', Pt(10), TEXT_GRAY)
    add_para(tf, '  RAG-Agent："根据校园生活指南，图书馆\n  开放时间为 8:00-22:00..."  <- 准确回答', Pt(10), GREEN)

    # BOTTOM comparison
    by = Inches(4.8)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), by, Inches(12.3), Inches(2.4))
    _fill(bar, CARD_BG); _set_line(bar, BORDER_COLOR, Pt(0.5))
    add_textbox(slide, Inches(0.8), Inches(4.95), Inches(11.5), Inches(0.3),
                "无RAG vs 有RAG -- 对比流程", Pt(14), PRIMARY_DARK, True)

    tf = add_rich_tf(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(1.5))
    add_para(tf, "X  无RAG流程：", Pt(11), RED, True)
    add_para(tf, "  用户提问 -> LLM直接回答 -> 可能产生幻觉", Pt(10), TEXT_DARK)
    add_para(tf, "  缺点：知识有限、无法溯源、可能编造信息", Pt(10), TEXT_GRAY)

    tf = add_rich_tf(slide, Inches(6.8), Inches(5.4), Inches(5.5), Inches(1.5))
    add_para(tf, "OK  RAG流程：", Pt(11), GREEN, True)
    add_para(tf, "  用户提问 -> 向量检索 -> 拼接到Prompt -> LLM生成", Pt(10), TEXT_DARK)
    add_para(tf, "  优势：知识实时更新、答案可溯源、准确率高", Pt(10), TEXT_GRAY)


# ============================================================
# SLIDE 3 - RAG PRINCIPLE
# ============================================================
def create_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "RAG 原理图",
                  "Retrieval -> Augmentation -> Generation  三步构建智能问答系统")

    col_w = Inches(3.5); col_h = Inches(3.5); gap = Inches(0.55)
    sx = Inches(0.55); sy = Inches(1.3)

    cols = [
        ("(1) 检索 Retrieval", PRIMARY_DARK,
         ["用户输入自然语言问题", "Embedding模型将问题转为向量",
          "在ChromaDB向量库中相似度搜索", "召回 Top-K 个最相关文档片段",
          "过滤低相似度结果 (threshold>0.3)"]),
        ("(2) 增强 Augmentation", PRIMARY,
         ["将检索到的文档片段拼接", "构建结构化Prompt模板",
          "注入系统指令 + 角色设定", "添加知识库上下文 {context}",
          "嵌入用户原始问题 {question}"]),
        ("(3) 生成 Generation", RGBColor(0x8A, 0x6A, 0xA6),
         ["将完整Prompt发送至DeepSeek", "大模型基于上下文推理回答",
          "生成结构化、可读的答案", "附带引用来源 (source标注)",
          "兜底策略：知识库外问题友好提示"]),
    ]

    for idx, (title, stripe_c, items) in enumerate(cols):
        cx = sx + idx * (col_w + gap)
        add_card(slide, cx, sy, col_w, col_h + Inches(0.3))
        add_rect(slide, cx, sy, col_w, Inches(0.06), stripe_c)
        add_textbox(slide, cx + Inches(0.2), sy + Inches(0.2),
                    col_w - Inches(0.4), Inches(0.4), title, Pt(15), stripe_c, True)
        tf = add_rich_tf(slide, cx + Inches(0.2), sy + Inches(0.7),
                         col_w - Inches(0.4), col_h - Inches(0.5))
        for item in items:
            add_para(tf, "  > " + item, Pt(10), TEXT_DARK, sb=Pt(3), sa=Pt(3))
        if idx < 2:
            add_gold_arrow(slide, cx + col_w + Inches(0.05),
                           sy + col_h / 2 - Inches(0.12))

    # Bottom example
    fy = Inches(5.3)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), fy, Inches(12.3), Inches(1.8))
    _fill(bar, CARD_BG); _set_line(bar, BORDER_COLOR, Pt(0.5))
    add_textbox(slide, Inches(0.8), fy + Inches(0.1), Inches(11.5), Inches(0.3),
                '示例流程：用户提问 "新生宿舍怎么分配？"', Pt(13), PRIMARY_DARK, True)
    tf = add_rich_tf(slide, Inches(0.8), fy + Inches(0.5), Inches(11.5), Inches(1.1))
    add_para(tf, 'Step 1 检索  ->  "新生宿舍..." -> Embedding -> ChromaDB查询 -> '
                  'Top-3文档片段 [校园生活指南.md, 新生入学指南.md, ...]', Pt(10), TEXT_DARK)
    add_para(tf, 'Step 2 增强  ->  将片段拼入Prompt模板：系统角色 + {context} + {question}',
             Pt(10), TEXT_DARK)
    add_para(tf, 'Step 3 生成  ->  DeepSeek返回："根据新生入学指南，宿舍由学校统一分配，按专业和班级..."',
             Pt(10), GREEN)


# ============================================================
# SLIDE 4 - TECH STACK
# ============================================================
def create_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "技术栈", "项目使用的核心技术组件与工具链")

    tools = [
        ("DeepSeek Chat", "LLM 大语言模型", "核心引擎",
         "作为核心推理引擎，基于DeepSeek大模型API提供高质量的中文理解和生成能力，支持32K上下文窗口，性价比优异。"),
        ("Qwen Embedding", "Embedding 向量化模型", "向量化",
         "将文档文本转换为高维向量表示（1024维），支持中文语义理解，确保检索精度的基础组件。采用text-embedding-v3模型。"),
        ("ChromaDB", "向量数据库", "存储检索",
         "轻量级开源向量数据库，存储文档Embedding向量并支持高效的余弦相似度检索，部署简单、查询速度快。"),
        ("Gradio", "Web UI 框架", "界面框架",
         "快速构建AI应用Web界面的Python库，支持ChatInterface组件，提供流式输出、多轮对话、文件上传等功能。"),
        ("Qoder CLI", "AI 辅助开发工具", "开发工具",
         "AI编程助手命令行工具，辅助代码生成、调试和项目管理，提升开发效率，在项目全生命周期中发挥重要作用。"),
    ]

    rh = Inches(0.78); rg = Inches(0.12); sy = Inches(1.2)
    lx = Inches(0.5); cw = Inches(12.3)

    for i, (name, role, badge, desc) in enumerate(tools):
        ry = sy + i * (rh + rg)
        add_card(slide, lx, ry, cw, rh)
        stripe_c = PURPLE_SHADES[i]
        add_rect(slide, lx, ry, Inches(0.07), rh, stripe_c)
        add_textbox(slide, lx + Inches(0.25), ry + Inches(0.08),
                    Inches(2.0), Inches(0.3), name, Pt(14), PRIMARY_DARK, True)
        add_badge(slide, lx + Inches(0.25), ry + Inches(0.4),
                  badge, stripe_c, TEXT_LIGHT, Pt(8))
        add_textbox(slide, lx + Inches(2.5), ry + Inches(0.12),
                    Inches(9.5), Inches(0.55), desc, Pt(10), TEXT_DARK)


# ============================================================
# SLIDE 5 - KNOWLEDGE BASE
# ============================================================
def create_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "知识库设计",
                  "10份校园主题Markdown文档，覆盖校园生活全场景")

    docs = [
        ("校园生活指南.md", "宿舍、食堂、图书馆等日常生活信息", "3,200"),
        ("教务与选课指南.md", "选课流程、学分要求、考试安排", "2,850"),
        ("奖助学金政策.md", "各类奖学金、助学金申请条件与流程", "2,100"),
        ("信息技术服务指南.md", "校园网、VPN、一卡通等IT服务", "1,950"),
        ("新生入学指南.md", "报到流程、军训安排、入学教育", "3,450"),
        ("学生社团活动指南.md", "社团介绍、活动申请、经费管理", "1,800"),
        ("安全与医疗指南.md", "校园安全、医疗就诊、心理咨询", "2,400"),
        ("交通与出行指南.md", "校车、公交、周边交通攻略", "1,650"),
        ("图书馆使用指南.md", "借阅规则、自习室、电子资源", "2,350"),
        ("毕业与就业指南.md", "毕业条件、就业指导、考研信息", "2,780"),
    ]

    tl = Inches(0.5); tt = Inches(1.2); tw = Inches(8.2)
    hh = Inches(0.35); rh = Inches(0.38)
    n_rows = len(docs) + 2
    ts = slide.shapes.add_table(n_rows, 3, tl, tt, tw, hh + rh * (n_rows - 1))
    tbl = ts.table
    tbl.columns[0].width = Inches(3.3)
    tbl.columns[1].width = Inches(3.5)
    tbl.columns[2].width = Inches(1.4)

    for ci, h in enumerate(["文档名称", "内容范围", "字数"]):
        fmt_cell(tbl.cell(0, ci), h, Pt(10), TEXT_LIGHT, True, PRIMARY_DARK, PP_ALIGN.CENTER)
    for ri, (name, scope, words) in enumerate(docs):
        bg = CARD_BG if ri % 2 == 0 else GRAY_BG
        fmt_cell(tbl.cell(ri+1, 0), "  " + name, Pt(9), TEXT_DARK, False, bg)
        fmt_cell(tbl.cell(ri+1, 1), scope, Pt(9), TEXT_DARK, False, bg)
        fmt_cell(tbl.cell(ri+1, 2), words, Pt(9), TEXT_DARK, False, bg, PP_ALIGN.CENTER)
    tr = len(docs) + 1
    fmt_cell(tbl.cell(tr, 0), "  合计 10 份文档", Pt(10), TEXT_LIGHT, True, PRIMARY_DARK)
    fmt_cell(tbl.cell(tr, 1), "覆盖 12 大类校园场景", Pt(10), TEXT_LIGHT, True, PRIMARY_DARK, PP_ALIGN.CENTER)
    fmt_cell(tbl.cell(tr, 2), "24,530", Pt(10), TEXT_LIGHT, True, PRIMARY_DARK, PP_ALIGN.CENTER)

    rx = Inches(9.1); rw = Inches(3.8)

    # Card 1: Chunking strategy
    add_card(slide, rx, Inches(1.2), rw, Inches(2.3))
    add_rect(slide, rx, Inches(1.2), rw, Inches(0.05), PRIMARY)
    tf = add_rich_tf(slide, rx + Inches(0.2), Inches(1.4), rw - Inches(0.4), Inches(1.9))
    add_para(tf, "切片策略", Pt(15), PRIMARY_DARK, True, sb=Pt(4))
    add_para(tf, "", Pt(4))
    add_para(tf, "  chunk_size  = 300", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  overlap      = 50", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  分隔符       = \\n\\n (段落级)", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  总计切片     ~ 160 个片段", Pt(11), TEXT_DARK, sb=Pt(3))

    # Card 2: Technical details
    add_card(slide, rx, Inches(3.8), rw, Inches(2.5))
    add_rect(slide, rx, Inches(3.8), rw, Inches(0.05), ACCENT)
    tf = add_rich_tf(slide, rx + Inches(0.2), Inches(4.0), rw - Inches(0.4), Inches(2.1))
    add_para(tf, "技术细节", Pt(15), ACCENT, True, sb=Pt(4))
    add_para(tf, "", Pt(4))
    add_para(tf, "  Embedding维度  = 1024", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  检索Top-K      = 5", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  相似度阈值      = 0.3", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  Batch Size      = 16", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  向量距离度量    = Cosine", Pt(11), TEXT_DARK, sb=Pt(3))
    add_para(tf, "  文档格式        = Markdown", Pt(11), TEXT_DARK, sb=Pt(3))

# ============================================================
# SLIDE 6 - CORE CODE
# ============================================================
def create_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "核心代码",
                  "RAG Agent Prompt模板 & 向量检索核心函数")

    code_w = Inches(5.9); code_h = Inches(2.75)
    code_y = Inches(1.15)

    # Left: Prompt template
    add_textbox(slide, Inches(0.5), Inches(1.08), Inches(5.9), Inches(0.2),
                "rag_agent.py -- Prompt模板", Pt(9), PRIMARY_DARK, True, font=CODE_FONT)
    prompt_code = (
        'SYSTEM_PROMPT = """你是楚雄师范学院校园智能助手。\n'
        '请根据以下知识库内容回答问题。\n\n'
        '知识库内容：\n'
        '{context}\n\n'
        '用户问题：{question}\n\n'
        '回答要求：\n'
        '1. 基于知识库内容回答\n'
        '2. 如果知识库没有相关信息，请明确说明\n'
        '3. 保持友好、专业的语气"""'
    )
    add_code_block(slide, Inches(0.5), code_y, code_w, code_h, prompt_code)

    # Right: search() function
    add_textbox(slide, Inches(6.9), Inches(1.08), Inches(5.9), Inches(0.2),
                "vector_store.py -- search()函数", Pt(9), PRIMARY_DARK, True, font=CODE_FONT)
    search_code = (
        'def search(self, query: str, top_k: int = 5):\n'
        '    """检索最相关的文档片段"""\n'
        '    query_embedding = self.embedding_model.encode(query)\n'
        '    \n'
        '    results = self.collection.query(\n'
        '        query_embeddings=[query_embedding],\n'
        '        n_results=top_k\n'
        '    )\n'
        '    \n'
        '    return [\n'
        '        {"content": doc, "metadata": meta, "distance": dist}\n'
        '        for doc, meta, dist in zip(\n'
        '            results["documents"][0],\n'
        '            results["metadatas"][0],\n'
        '            results["distances"][0]\n'
        '        )\n'
        '    ]'
    )
    add_code_block(slide, Inches(6.9), code_y, code_w, code_h, search_code)

    # File structure tree
    fs_y = Inches(4.15)
    add_textbox(slide, Inches(0.5), Inches(4.08), Inches(12), Inches(0.2),
                "project/  项目文件结构", Pt(9), PRIMARY_DARK, True, font=CODE_FONT)
    fs_code = (
        'project/\n'
        '  rag_agent.py          # RAG 核心逻辑\n'
        '  vector_store.py       # 向量数据库封装\n'
        '  document_loader.py    # 文档加载与切片\n'
        '  web_app.py            # Gradio Web 界面\n'
        '  config.py             # 配置文件\n'
        '  requirements.txt      # 依赖列表\n'
        '  data/\n'
        '    knowledge_base/     # 知识库文档'
    )
    add_code_block(slide, Inches(0.5), fs_y, Inches(12.3), Inches(1.75), fs_code)

    # Bottom stats
    stats_y = Inches(6.15)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), stats_y, Inches(12.3), Inches(1.1))
    _fill(bar, PRIMARY_DARK); _no_line(bar)
    add_textbox(slide, Inches(0.8), stats_y + Inches(0.1), Inches(11.5), Inches(0.3),
                "项目统计", Pt(14), ACCENT, True)
    add_textbox(slide, Inches(0.8), stats_y + Inches(0.55), Inches(11.5), Inches(0.35),
                "Python文件 5个  |  代码行数 ~800行  |  知识库文档 10份  |  切片片段 ~160个  |  测试场景 42个",
                Pt(12), TEXT_LIGHT)


# ============================================================
# SLIDE 7 - WEB INTERFACE
# ============================================================
def create_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "Web 界面展示",
                  "基于Gradio构建的校园智能问答Web应用")

    # LEFT - UI features
    add_card(slide, Inches(0.5), Inches(1.2), Inches(5.8), Inches(4.5))
    tf = add_rich_tf(slide, Inches(0.8), Inches(1.4), Inches(5.2), Inches(4.0))
    add_para(tf, "UI 功能特性", Pt(16), PRIMARY_DARK, True, sb=Pt(6))
    add_para(tf, "", Pt(4))
    features = [
        ("[1] ChatInterface 对话组件", "流式输出、Markdown渲染"),
        ("[2] 多轮对话记忆", "自动维护对话历史上下文"),
        ("[3] 参考来源展示", "每条回答附文档来源链接"),
        ("[4] 知识库管理面板", "支持文档上传与索引重建"),
        ("[5] 模型参数调节", "Temperature、Top-K等可调"),
        ("[6] 示例问题推荐", "预设快捷提问按钮"),
        ("[7] 一键清空对话", "重置对话历史"),
        ("[8] 响应式布局", "适配不同屏幕尺寸"),
    ]
    for name, desc in features:
        add_para(tf, name, Pt(12), TEXT_DARK, True, sb=Pt(5), sa=Pt(1))
        add_para(tf, "      " + desc, Pt(10), TEXT_GRAY, sb=Pt(1), sa=Pt(3))

    # RIGHT - Screenshot placeholder
    add_card(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(4.5),
             RGBColor(0xEE, 0xEA, 0xF4))
    add_rect(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(0.05), PRIMARY)
    tf = add_rich_tf(slide, Inches(7.1), Inches(2.8), Inches(5.2), Inches(1.5))
    add_para(tf, "Web界面截图区域", Pt(20), PRIMARY_LIGHT, True, PP_ALIGN.CENTER)
    add_para(tf, "", Pt(8))
    add_para(tf, "请在此处插入Gradio Web界面截图", Pt(12), TEXT_GRAY, False, PP_ALIGN.CENTER)
    add_para(tf, "展示智能问答交互效果", Pt(11), TEXT_GRAY, False, PP_ALIGN.CENTER)

    # BOTTOM highlight
    hy = Inches(6.0)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), hy, Inches(12.3), Inches(1.15))
    _fill(bar, PRIMARY_DARK); _no_line(bar)
    add_textbox(slide, Inches(0.8), hy + Inches(0.1), Inches(11.5), Inches(0.3),
                "关键技术亮点", Pt(13), ACCENT, True)
    add_textbox(slide, Inches(0.8), hy + Inches(0.5), Inches(11.5), Inches(0.4),
                "多轮对话上下文管理  |  检索结果来源标注与引用  |  兜底策略保障回答质量  |  Markdown格式化输出",
                Pt(11), TEXT_LIGHT)


# ============================================================
# SLIDE 8 - TEST RESULTS
# ============================================================
def create_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "测试结果",
                  "42个校园场景问答测试，覆盖12大类，通过率100%")

    # Main pass rate card
    add_card(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.1))
    add_rect(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.06), GREEN)
    tf = add_rich_tf(slide, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.8))
    add_para(tf, "42个问答场景全部通过  通过率 100%", Pt(22), GREEN, True, PP_ALIGN.CENTER,
             sb=Pt(8), sa=Pt(4))
    add_para(tf, "测试日期：2025年3月 | 模型版本：DeepSeek-V3 | 知识库版本：v1.0",
             Pt(11), TEXT_GRAY, False, PP_ALIGN.CENTER)

    # Test categories grid
    categories = [
        ("宿舍生活", "4/4"), ("食堂餐饮", "3/3"), ("图书馆", "3/3"),
        ("奖助学金", "3/3"), ("教务选课", "4/4"), ("考试毕业", "3/3"),
        ("信息技术", "3/3"), ("交通出行", "3/3"), ("社团活动", "3/3"),
        ("安全医疗", "3/3"), ("新生入学", "4/4"), ("毕业就业", "3/3"),
    ]

    grid_x = Inches(0.5); grid_y = Inches(2.55)
    cell_w = Inches(1.95); cell_h = Inches(1.0)
    gap_x = Inches(0.08); gap_y = Inches(0.08)
    cols_per_row = 6

    for i, (name, score) in enumerate(categories):
        row = i // cols_per_row; col = i % cols_per_row
        cx = grid_x + col * (cell_w + gap_x)
        cy = grid_y + row * (cell_h + gap_y)
        add_card(slide, cx, cy, cell_w, cell_h)
        add_rect(slide, cx, cy, cell_w, Inches(0.04), GREEN)
        add_textbox(slide, cx + Inches(0.1), cy + Inches(0.08), cell_w - Inches(0.2), Inches(0.3),
                    name, Pt(12), PRIMARY_DARK, True, PP_ALIGN.CENTER)
        add_textbox(slide, cx + Inches(0.1), cy + Inches(0.48), cell_w - Inches(0.2), Inches(0.3),
                    "OK  " + score, Pt(11), GREEN, True, PP_ALIGN.CENTER)

    # Bottom: fallback strategy
    fy = Inches(4.85)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), fy, Inches(12.3), Inches(2.3))
    _fill(bar, CARD_BG); _set_line(bar, BORDER_COLOR, Pt(0.5))
    add_textbox(slide, Inches(0.8), fy + Inches(0.1), Inches(11.5), Inches(0.3),
                "兜底策略验证", Pt(15), PRIMARY_DARK, True)
    tf = add_rich_tf(slide, Inches(0.8), fy + Inches(0.55), Inches(11.5), Inches(1.5))
    add_para(tf, "测试场景：知识库外问题（如：\"今天天气怎么样？\"）", Pt(12), TEXT_DARK, True, sb=Pt(4))
    add_para(tf, "  - 检索结果相似度均低于0.3阈值", Pt(11), TEXT_DARK, sb=Pt(2))
    add_para(tf, "  - 系统正确触发兜底策略", Pt(11), TEXT_DARK, sb=Pt(2))
    add_para(tf, '  - 返回友好提示："抱歉，我暂时无法回答这个问题，建议您联系学校相关部门获取帮助。"',
             Pt(11), GREEN, sb=Pt(2))
    add_para(tf, "", Pt(4))
    add_para(tf, "结论：兜底策略有效运作，无幻觉输出，用户体验得到保障。", Pt(12), PRIMARY_DARK, True, sb=Pt(4))


# ============================================================
# SLIDE 9 - DEBUG CASES (2x2 bug cards)
# ============================================================
def create_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "问题与解决",
                  "开发过程中遇到的典型问题及解决方案")

    bugs = [
        {
            "name": "split_text死循环",
            "symptom": "长文档切片时程序卡死，CPU占用100%",
            "cause": "递归分割逻辑未处理单字符超长情况，分隔符split后list为空导致无限递归",
            "fix": "添加最大递归深度限制(max_depth=10)，并在chunk_size小于字符串长度时强制截断",
        },
        {
            "name": "Qwen Embedding批量超限",
            "symptom": "批量Embedding请求返回HTTP 413错误",
            "cause": "Qwen Embedding API单次请求token上限为2048，批量发送16个chunk时超限",
            "fix": "将batch_size从16调整为8，并添加自动重试机制(最多3次)，失败时逐条处理",
        },
        {
            "name": "Gradio 6.x API不兼容",
            "symptom": "ChatInterface组件的respond方法签名错误",
            "cause": "Gradio从5.x升级到6.x后，chat函数签名从(message, history)变为(message, history, request)",
            "fix": "更新函数签名添加request参数，并在requirements.txt中锁定gradio>=6.0版本",
        },
        {
            "name": "文档来源未显示",
            "symptom": "回答中没有显示引用的文档来源",
            "cause": "检索返回的metadata中source字段未正确传递到Prompt的context中",
            "fix": "在构建context时显式拼接source信息，格式为[来源: 文件名.md]，确保LLM能引用",
        },
    ]

    card_w = Inches(5.8); card_h = Inches(2.7)
    positions = [
        (Inches(0.5), Inches(1.15)),
        (Inches(6.8), Inches(1.15)),
        (Inches(0.5), Inches(4.15)),
        (Inches(6.8), Inches(4.15)),
    ]

    for i, bug in enumerate(bugs):
        bx, by = positions[i]
        add_card(slide, bx, by, card_w, card_h)
        add_rect(slide, bx, by, card_w, Inches(0.05), RED)
        add_badge(slide, bx + Inches(0.2), by + Inches(0.15),
                  "Bug #" + str(i+1) + "  " + bug["name"], RED, TEXT_LIGHT, Pt(9))

        tf = add_rich_tf(slide, bx + Inches(0.2), by + Inches(0.65),
                         card_w - Inches(0.4), card_h - Inches(0.75))
        add_para(tf, "症状: " + bug["symptom"], Pt(10), RED, True, sb=Pt(3), sa=Pt(2))
        add_para(tf, "原因: " + bug["cause"], Pt(10), TEXT_DARK, False, sb=Pt(2), sa=Pt(2))
        add_para(tf, "修复: " + bug["fix"], Pt(10), GREEN, False, sb=Pt(2), sa=Pt(2))


# ============================================================
# SLIDE 10 - SUMMARY
# ============================================================
def create_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, BG_COLOR)
    add_title_bar(slide, "技术复盘 & 拓展学习路径",
                  "项目总结、技能收获与未来方向")

    # LEFT - Skills table
    add_textbox(slide, Inches(0.5), Inches(1.08), Inches(6), Inches(0.3),
                "技能收获", Pt(16), PRIMARY_DARK, True)

    skills = [
        ("大模型API调用", "掌握DeepSeek等大模型API的调用与参数调优"),
        ("RAG原理与实践", "深入理解检索增强生成的完整技术链路"),
        ("Prompt Engineering", "学会设计结构化Prompt模板提升回答质量"),
        ("向量数据库", "使用ChromaDB进行文档Embedding存储与检索"),
        ("文档处理", "掌握Markdown解析、文本切片与预处理技术"),
        ("AI辅助开发", "利用Qoder CLI等AI工具提升开发效率"),
        ("Web应用开发", "使用Gradio快速构建AI应用Web界面"),
    ]

    tl2 = Inches(0.5); tt2 = Inches(1.5); tw2 = Inches(6.5)
    ts = slide.shapes.add_table(len(skills) + 1, 2, tl2, tt2, tw2,
                                Inches(0.38) * (len(skills) + 1))
    tbl = ts.table
    tbl.columns[0].width = Inches(2.2); tbl.columns[1].width = Inches(4.3)

    fmt_cell(tbl.cell(0, 0), "  技能领域", Pt(10), TEXT_LIGHT, True, PRIMARY_DARK)
    fmt_cell(tbl.cell(0, 1), "  具体收获", Pt(10), TEXT_LIGHT, True, PRIMARY_DARK)

    for ri, (skill, detail) in enumerate(skills):
        bg = CARD_BG if ri % 2 == 0 else GRAY_BG
        fmt_cell(tbl.cell(ri+1, 0), "  " + skill, Pt(10), PRIMARY_DARK, True, bg)
        fmt_cell(tbl.cell(ri+1, 1), detail, Pt(9), TEXT_DARK, False, bg)

    # RIGHT - Learning paths
    add_textbox(slide, Inches(7.5), Inches(1.08), Inches(5.5), Inches(0.3),
                "拓展学习路径", Pt(16), PRIMARY_DARK, True)

    paths = [
        ("(1) 框架化开发", "将RAG系统重构为LangChain/LlamaIndex框架，提升可维护性和扩展性"),
        ("(2) 多模态RAG", "支持图片、PDF、表格等多模态文档的检索与问答"),
        ("(3) Agentic RAG", "引入Agent自主决策能力，动态选择检索策略与工具调用"),
        ("(4) Graph RAG", "构建知识图谱增强检索，提升实体关系推理能力"),
        ("(5) 评估与优化", "引入RAGAS等评估框架，量化检索与生成质量，持续迭代优化"),
    ]

    for i, (title, desc) in enumerate(paths):
        py = Inches(1.55) + i * Inches(1.05)
        add_card(slide, Inches(7.5), py, Inches(5.5), Inches(0.95))
        add_rect(slide, Inches(7.5), py, Inches(0.07), Inches(0.95), PURPLE_SHADES[i])
        add_textbox(slide, Inches(7.75), py + Inches(0.08), Inches(5.0), Inches(0.3),
                    title, Pt(13), PRIMARY_DARK, True)
        add_textbox(slide, Inches(7.75), py + Inches(0.42), Inches(5.0), Inches(0.45),
                    desc, Pt(10), TEXT_DARK)

    # Bottom thank-you
    ty_y = Inches(6.75)
    add_rect(slide, Inches(0), ty_y, SLIDE_W, Inches(0.75), PRIMARY_DARK)
    add_textbox(slide, Inches(0.5), ty_y + Inches(0.12), Inches(12.3), Inches(0.5),
                "谢谢！欢迎提问与交流",
                Pt(22), ACCENT, True, PP_ALIGN.CENTER)


# ============================================================
# MAIN
# ============================================================
def main():
    print("Creating presentation...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    create_slide_1(prs)
    create_slide_2(prs)
    create_slide_3(prs)
    create_slide_4(prs)
    create_slide_5(prs)
    create_slide_6(prs)
    create_slide_7(prs)
    create_slide_8(prs)
    create_slide_9(prs)
    create_slide_10(prs)

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
